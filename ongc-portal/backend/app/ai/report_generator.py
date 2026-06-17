import os
import tempfile
from datetime import datetime
from typing import Optional
from app.ai.llm_client import llm
from app.config import settings

class ReportGenerator:
    async def generate_report(self, topic: str, format: str = "pdf", user_name: str = "User") -> dict:
        prompt = f"""Generate a professional business report about: {topic}

Include:
1. Executive Summary
2. Key Findings and Analysis
3. Supporting Data Points
4. Recommendations
5. Next Steps

Format the report with clear sections and professional language."""
        content = await llm.generate(prompt, "You are a business report writer for ONGC. Produce comprehensive, professional reports.")

        if format == "pdf":
            path = await self._generate_pdf(topic, content, user_name)
        elif format == "docx":
            path = await self._generate_docx(topic, content, user_name)
        elif format == "pptx":
            path = await self._generate_pptx(topic, content, user_name)
        else:
            path = None

        return {
            "topic": topic,
            "format": format,
            "content": content,
            "file_path": path,
            "generated_at": datetime.utcnow().isoformat(),
        }

    async def _generate_pdf(self, topic: str, content: str, user_name: str) -> str:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak

        fd, path = tempfile.mkstemp(suffix=".pdf", prefix=f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}_")
        os.close(fd)

        doc = SimpleDocTemplate(path, pagesize=A4,
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=72)

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle', parent=styles['Title'],
            fontSize=20, spaceAfter=20, textColor=colors.HexColor("#0b3d91"),
        )
        heading_style = ParagraphStyle(
            'CustomHeading', parent=styles['Heading2'],
            fontSize=14, spaceAfter=10, textColor=colors.HexColor("#0b3d91"),
        )
        body_style = ParagraphStyle(
            'CustomBody', parent=styles['Normal'],
            fontSize=10, spaceAfter=8, leading=14,
        )

        story = []
        story.append(Paragraph(f"Data Vision", styles['Normal']))
        story.append(Spacer(1, 6))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Spacer(1, 6))
        story.append(Paragraph(f"By: {user_name}", styles['Normal']))
        story.append(Spacer(1, 20))

        lines = content.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
            elif line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or \
                 line.startswith("4.") or line.startswith("5.") or \
                 any(line.startswith(h) for h in ["Executive Summary", "Key Findings", "Analysis",
                                                    "Recommendations", "Next Steps", "Supporting Data"]):
                story.append(Paragraph(f"<b>{line}</b>", heading_style))
            else:
                story.append(Paragraph(line, body_style))

        doc.build(story)
        return path

    async def _generate_docx(self, topic: str, content: str, user_name: str) -> str:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)

        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run("Data Vision")
        run.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(11, 61, 145)

        doc.add_paragraph(f"Report: {topic}")
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_paragraph(f"By: {user_name}")
        doc.add_paragraph()

        lines = content.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if any(line.startswith(h) for h in ["Executive Summary", "Key Findings", "Analysis",
                                                  "Recommendations", "Next Steps", "Supporting Data"]):
                doc.add_heading(line, level=1)
            elif line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or \
                 line.startswith("4.") or line.startswith("5."):
                doc.add_paragraph(line, style='List Number')
            else:
                doc.add_paragraph(line)

        fd, path = tempfile.mkstemp(suffix=".docx", prefix=f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}_")
        os.close(fd)
        doc.save(path)
        return path

    async def _generate_pptx(self, topic: str, content: str, user_name: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 61, 145)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = topic
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"ONGC - {datetime.now().strftime('%B %Y')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        p2.alignment = PP_ALIGN.CENTER

        sections = content.split("\n\n")
        for section in sections:
            lines = section.strip().split("\n")
            if not lines or not lines[0].strip():
                continue
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = lines[0].strip()[:100]

            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
            tf = txBox.text_frame
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.space_after = Pt(6)

        fd, path = tempfile.mkstemp(suffix=".pptx", prefix=f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}_")
        os.close(fd)
        prs.save(path)
        return path

    async def list_reports(self, report_dir: str = None) -> list[dict]:
        import glob
        if not report_dir:
            report_dir = tempfile.gettempdir()
        reports = []
        for ext in ["*.pdf", "*.docx", "*.pptx"]:
            pattern = os.path.join(report_dir, f"report_*{ext}")
            for fpath in glob.glob(pattern):
                stats = os.stat(fpath)
                reports.append({
                    "path": fpath,
                    "filename": os.path.basename(fpath),
                    "size": stats.st_size,
                    "modified": datetime.fromtimestamp(stats.st_mtime).isoformat(),
                })
        return sorted(reports, key=lambda x: x["modified"], reverse=True)[:20]

report_gen = ReportGenerator()
