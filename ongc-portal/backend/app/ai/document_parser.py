import os
import csv
import io
from typing import List, Tuple
from app.config import settings

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".txt"}

def parse_file(file_path: str) -> Tuple[str, List[dict]]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf_marker(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".txt":
        return parse_txt(file_path)
    return "", []

def parse_pdf_marker(file_path: str) -> Tuple[str, List[dict]]:
    if settings.MARKER_ENABLED:
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            converter = PdfConverter(artifact_dict=create_model_dict())
            rendered = converter(file_path)
            full_text = rendered.markdown
            pages = [{"page": i + 1, "text": p.markdown or ""} for i, p in enumerate(rendered.pages)]
            if full_text.strip():
                return full_text, pages
        except Exception:
            pass
    return parse_pdf_paddle(file_path)

def parse_pdf_paddle(file_path: str) -> Tuple[str, List[dict]]:
    if settings.PADDLE_ENABLED:
        try:
            from app.ai.paddle_ocr_service import parse_pdf_paddle_ppocrv5
            text, pages = parse_pdf_paddle_ppocrv5(file_path)
            if text and text.strip():
                return text, pages or []
        except Exception:
            pass
    return parse_pdf_ocr(file_path)

def parse_pdf_ocr(file_path: str) -> Tuple[str, List[dict]]:
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(file_path, dpi=300)
        full_text = ""
        pages = []
        for i, img in enumerate(images):
            page_text = pytesseract.image_to_string(img, lang="eng")
            full_text += page_text + "\n"
            pages.append({"page": i + 1, "text": page_text})
        if full_text.strip():
            return full_text, pages
    except Exception:
        pass
    return parse_pdf_pypdf2(file_path)

def parse_pdf_pypdf2(file_path: str, marker_error: str = None) -> Tuple[str, List[dict]]:
    import PyPDF2
    text = ""
    pages = []
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                pages.append({"page": i + 1, "text": page_text})
        header = ""
        if marker_error:
            header = f"[Marker unavailable: {marker_error}. Falling back to PyPDF2.]\n"
        return header + text, pages
    except Exception as e:
        text = f"[Error parsing PDF: {str(e)}]"
    return text, pages

def parse_docx(file_path: str) -> Tuple[str, List[dict]]:
    import docx
    text = ""
    paragraphs = []
    try:
        doc = docx.Document(file_path)
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text += para.text + "\n"
                paragraphs.append({"paragraph": i + 1, "text": para.text})
    except Exception as e:
        text = f"[Error parsing DOCX: {str(e)}]"
    return text, paragraphs

def parse_xlsx(file_path: str) -> Tuple[str, List[dict]]:
    import openpyxl
    text = ""
    sheets = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = f"Sheet: {sheet_name}\n"
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join([str(c) if c is not None else "" for c in row])
                sheet_text += row_text + "\n"
                rows.append(list(row))
            text += sheet_text + "\n"
            sheets.append({"sheet": sheet_name, "rows": rows})
        wb.close()
    except Exception as e:
        text = f"[Error parsing XLSX: {str(e)}]"
    return text, sheets

def parse_csv(file_path: str) -> Tuple[str, List[dict]]:
    text = ""
    rows_list = []
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                row_text = "\t".join(row)
                text += row_text + "\n"
                rows_list.append({"row": i + 1, "text": row_text})
    except Exception as e:
        text = f"[Error parsing CSV: {str(e)}]"
    return text, rows_list

def parse_pptx(file_path: str) -> Tuple[str, List[dict]]:
    from pptx import Presentation
    text = ""
    slides = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                    slide_parts.append(shape.text)
            if slide_text:
                text += f"Slide {i + 1}:\n{slide_text}\n"
                slides.append({"slide": i + 1, "text": slide_text, "parts": slide_parts})
    except Exception as e:
        text = f"[Error parsing PPTX: {str(e)}]"
    return text, slides

def parse_txt(file_path: str) -> Tuple[str, List[dict]]:
    text = ""
    lines = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            for i, line in enumerate(f):
                text += line
                if line.strip():
                    lines.append({"line": i + 1, "text": line.rstrip()})
    except Exception as e:
        text = f"[Error parsing TXT: {str(e)}]"
    return text, lines

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        i += chunk_size - overlap
        if i >= len(words):
            break
    return chunks

def chunk_text_with_pages(text: str, pages: List[dict], chunk_size: int = 500, overlap: int = 50) -> List[dict]:
    full_text = text
    words = full_text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk_words = words[i:i + chunk_size]
        chunk_text_str = " ".join(chunk_words)
        if not chunk_text_str.strip():
            i += chunk_size - overlap
            continue
        start_char_idx = len(" ".join(words[:i]))
        end_char_idx = start_char_idx + len(chunk_text_str)
        page_refs = []
        for p in pages:
            p_start = full_text.find(p["text"])
            if p_start == -1:
                continue
            p_end = p_start + len(p["text"])
            if p_start < end_char_idx and p_end > start_char_idx:
                page_refs.append(p.get("page", 0))
        chunks.append({
            "text": chunk_text_str,
            "chunk_index": len(chunks),
            "page_numbers": list(set(page_refs)) if page_refs else [],
        })
        i += chunk_size - overlap
    return chunks
